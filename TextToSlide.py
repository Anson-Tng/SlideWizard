# Name      :   Anson Ting
# File      :   TextToSlide.py
# Purposes  :   Convert all the generated text to txt file.

import copy
from inspect import Parameter
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import math
import logging


def _copy_slide(pres, source):
    # Copy slide feature,
    # Credit: https://github.com/scanny/python-pptx/issues/132#issuecomment-1098398796
    dest = pres.slides.add_slide(source.slide_layout)
    for shape in dest.shapes:
        shape.element.getparent().remove(shape.element)

    for shape in source.shapes:
        new_shape = copy.deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

    for rel in source.part.rels:
        target = rel._target

        if "notesSlide" in rel.reltype:
            continue

    return dest


def _delete_slide(prs, slide):
    # Delete slide feature
    # Credit: https://github.com/scanny/python-pptx/issues/67#issuecomment-296135015

    # Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId]
               for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


def _add_image(slide, placeholder_id, image_url):
    # Add image feature for picture + text slide
    placeholder = slide.placeholders[placeholder_id]

    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size

    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width

    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)

    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # Placeholder width too wid
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 2
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 2
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side

    placeholder.width = 3000000
    placeholder.height = 3000000
    #placeholder.width = 5000000
    #placeholder.height = 5550000
    #placeholder.left = Inches(0.5)
    #placeholder.top = Inches(0.65)
    placeholder.left = Inches(2)
    placeholder.top = Inches(2)


def _find_placeholders(slide):
    # Find the placeholder id
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


# hardcoded, it uses one template powerpoint slide that contains 5 different pages.
pr = Presentation('Template.pptx')
# Once it works with multiple template, this can removed and all slide functions will have one more parameter: pr

# Template slide index
# 0 Intro
# 1 Picture - ToBeImplement
# 2 Bullet points
# 3 Text
# 4 Summary


def slide_intro(header, description=''):
    # Intro slide
    # Parameter : Title, Description(optional)
    intro = pr.slides[0]

    # Initial placeholder
    title = intro.shapes.title
    content = intro.placeholders[1]

    # Define placeholder content
    title.text = header
    content.text = description


def slide_with_pic(img, text):
    # Picture and text slide
    # Parameter : picture file, contents

    #img = "teamLogo.png"
    contentPage = _copy_slide(pr, pr.slides[1])
    _add_image(contentPage, 1, img)

    content = contentPage.placeholders[2]
    content.text = text


def slide_with_bullet_point(header, text):
    contentPage = _copy_slide(pr, pr.slides[2])
    title = contentPage.shapes.title
    title.text = header
    bpBox = contentPage.shapes.placeholders[1]
    
    # Clear any existing text in the bullet point box
    bpBox.text_frame.clear()

    # Split the text into main points and subpoints
    points = text.split(' - ')

    for point in points:
        if not point.strip():
            continue

        # Split subpoints from main points
        sub_points = point.split(' -- ')

        # Add the main point, stripping out any leading hyphens
        main_point_text = sub_points[0].lstrip('- ').strip()
        main_bullet = bpBox.text_frame.add_paragraph()
        main_bullet.text = main_point_text
        main_bullet.level = 0

        # Add any subpoints, ensuring they also have no leading hyphens
        for sub_point in sub_points[1:]:
            sub_point_text = sub_point.lstrip('- ').strip()
            sub_bullet = bpBox.text_frame.add_paragraph()
            sub_bullet.text = sub_point_text
            sub_bullet.level = 1  # This will indent the subpoint to create a hierarchy

def slide_with_text(header, text):
    # Header and text only slide
    # Parameter : Title, Description
    contentPage = _copy_slide(pr, pr.slides[3])
    title = contentPage.shapes.title
    content = contentPage.placeholders[1]

    title.text = header
    content.text = text


def slide_thank_you(header, text):
    #Just thank you!
    lastPage = _copy_slide(pr, pr.slides[4])


def slide_save(fileName):
    # Save feature
    # Hard coded
    # Delete all the template page
    _delete_slide(pr, pr.slides[1])
    _delete_slide(pr, pr.slides[1])
    _delete_slide(pr, pr.slides[1])
    _delete_slide(pr, pr.slides[1])

    # Save to file
    pr.save(fileName)



